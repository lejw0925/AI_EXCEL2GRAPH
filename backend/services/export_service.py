import logging
import base64
import json
from typing import Tuple, Dict, Any
from io import BytesIO
import asyncio

logger = logging.getLogger(__name__)

class ExportService:
    """导出服务"""
    
    def __init__(self):
        self.supported_formats = ['png', 'svg', 'pdf', 'markdown', 'iframe', 'pptx']
    
    async def export_text(self, chart_config: Dict, format_type: str, options: Any) -> str:
        """
        导出文本格式（markdown、iframe）
        
        Args:
            chart_config: 图表配置
            format_type: 导出格式
            options: 导出选项
            
        Returns:
            文本内容
        """
        try:
            if format_type == 'markdown':
                return await self._export_markdown(chart_config, options)
            elif format_type == 'iframe':
                return await self._export_iframe(chart_config, options)
            else:
                raise ValueError(f"不支持的文本格式: {format_type}")
                
        except Exception as e:
            logger.error(f"导出文本格式失败: {str(e)}")
            raise
    
    async def export_binary(self, chart_config: Dict, format_type: str, options: Any) -> Tuple[bytes, str]:
        """
        导出二进制格式（png、svg、pdf、pptx）
        
        Args:
            chart_config: 图表配置
            format_type: 导出格式
            options: 导出选项
            
        Returns:
            (文件内容, MIME类型)
        """
        try:
            if format_type == 'png':
                return await self._export_png(chart_config, options)
            elif format_type == 'svg':
                return await self._export_svg(chart_config, options)
            elif format_type == 'pdf':
                return await self._export_pdf(chart_config, options)
            elif format_type == 'pptx':
                return await self._export_pptx(chart_config, options)
            else:
                raise ValueError(f"不支持的二进制格式: {format_type}")
                
        except Exception as e:
            logger.error(f"导出二进制格式失败: {str(e)}")
            raise
    
    async def _export_markdown(self, chart_config: Dict, options: Any) -> str:
        """导出Markdown格式"""
        # 简化实现：生成包含配置的Markdown
        title = chart_config.get('title', {}).get('text', '图表')
        config_json = json.dumps(chart_config, ensure_ascii=False, indent=2)
        
        markdown_content = f"""# {title}

## 图表配置

```json
{config_json}
```

## 使用说明

此图表使用 ECharts 生成，可以直接在支持 ECharts 的环境中使用上述配置。

生成时间: {options.get('timestamp', '未知')}
图表类型: {chart_config.get('series', [{}])[0].get('type', '未知')}
"""
        
        return markdown_content
    
    async def _export_iframe(self, chart_config: Dict, options: Any) -> str:
        """导出iframe嵌入代码"""
        width = getattr(options, 'width', 800)
        height = getattr(options, 'height', 600)
        
        config_json = json.dumps(chart_config, ensure_ascii=False)
        
        iframe_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>ECharts 图表</title>
    <script src="https://cdn.jsdelivr.net/npm/echarts@5.4.3/dist/echarts.min.js"></script>
</head>
<body>
    <div id="chart" style="width: 100%; height: 100%;"></div>
    <script>
        var chartDom = document.getElementById('chart');
        var myChart = echarts.init(chartDom);
        var option = {config_json};
        myChart.setOption(option);
        
        // 响应式调整
        window.addEventListener('resize', function() {{
            myChart.resize();
        }});
    </script>
</body>
</html>"""
        
        iframe_code = f"""<iframe 
    src="data:text/html;charset=utf-8,{iframe_html.replace('"', '&quot;')}"
    width="{width}" 
    height="{height}" 
    frameborder="0">
</iframe>"""
        
        return iframe_code
    
    async def _export_png(self, chart_config: Dict, options: Any) -> Tuple[bytes, str]:
        """导出PNG格式"""
        # 这里使用简化实现，实际应该使用headless浏览器或服务端渲染
        # 由于没有安装相关依赖，返回模拟的PNG数据
        
        width = getattr(options, 'width', 800)
        height = getattr(options, 'height', 600)
        dpi = getattr(options, 'dpi', 300)
        
        # 创建一个简单的PNG头部（实际应该渲染真实图表）
        png_data = self._create_placeholder_image(width, height, 'PNG图表')
        
        return png_data, 'image/png'
    
    async def _export_svg(self, chart_config: Dict, options: Any) -> Tuple[bytes, str]:
        """导出SVG格式"""
        width = getattr(options, 'width', 800)
        height = getattr(options, 'height', 600)
        title = chart_config.get('title', {}).get('text', '图表')
        
        # 生成简单的SVG（实际应该渲染ECharts SVG）
        svg_content = f"""<?xml version="1.0" encoding="UTF-8"?>
<svg width="{width}" height="{height}" xmlns="http://www.w3.org/2000/svg">
    <rect width="100%" height="100%" fill="#f8f9fa"/>
    <text x="50%" y="50%" text-anchor="middle" font-family="Arial" font-size="24" fill="#333">
        {title or 'ECharts 图表'}
    </text>
    <text x="50%" y="60%" text-anchor="middle" font-family="Arial" font-size="12" fill="#666">
        图表类型: {chart_config.get('series', [{}])[0].get('type', '未知')}
    </text>
</svg>"""
        
        return svg_content.encode('utf-8'), 'image/svg+xml'
    
    async def _export_pdf(self, chart_config: Dict, options: Any) -> Tuple[bytes, str]:
        """导出PDF格式"""
        # 简化实现，实际应该使用报表库生成PDF
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            buffer = BytesIO()
            p = canvas.Canvas(buffer, pagesize=letter)
            
            title = chart_config.get('title', {}).get('text', '图表')
            p.drawString(100, 750, f"标题: {title}")
            p.drawString(100, 730, f"图表类型: {chart_config.get('series', [{}])[0].get('type', '未知')}")
            p.drawString(100, 710, "这是一个示例PDF，实际应包含渲染的图表")
            
            p.showPage()
            p.save()
            
            pdf_data = buffer.getvalue()
            buffer.close()
            
            return pdf_data, 'application/pdf'
            
        except ImportError:
            # 如果没有安装reportlab，返回简单的PDF内容描述
            simple_pdf = f"""PDF内容:
标题: {chart_config.get('title', {}).get('text', '图表')}
图表类型: {chart_config.get('series', [{}])[0].get('type', '未知')}
导出时间: {options.get('timestamp', '未知')}

注意: 此为简化版本，完整版本需要安装reportlab库"""
            
            return simple_pdf.encode('utf-8'), 'application/pdf'
    
    async def _export_pptx(self, chart_config: Dict, options: Any) -> Tuple[bytes, str]:
        """导出PowerPoint格式"""
        # 简化实现，实际应该使用python-pptx
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation()
            slide_layout = prs.slide_layouts[1]  # 标题和内容布局
            slide = prs.slides.add_slide(slide_layout)
            
            title = chart_config.get('title', {}).get('text', '图表')
            slide.shapes.title.text = title
            
            content = slide.placeholders[1]
            content.text = f"""图表信息:
• 类型: {chart_config.get('series', [{}])[0].get('type', '未知')}
• 系列数量: {len(chart_config.get('series', []))}
• 配置复杂度: {len(str(chart_config))} 字符

注意: 此为简化版本，完整版本应包含实际的图表图像"""
            
            buffer = BytesIO()
            prs.save(buffer)
            pptx_data = buffer.getvalue()
            buffer.close()
            
            return pptx_data, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            
        except ImportError:
            # 如果没有安装python-pptx，返回简单内容
            simple_pptx = f"""PowerPoint内容:
标题: {chart_config.get('title', {}).get('text', '图表')}
图表类型: {chart_config.get('series', [{}])[0].get('type', '未知')}

注意: 此为简化版本，完整版本需要安装python-pptx库"""
            
            return simple_pptx.encode('utf-8'), 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    
    def _create_placeholder_image(self, width: int, height: int, text: str) -> bytes:
        """创建占位图片"""
        # 简化实现：返回PNG文件头和基本数据
        # 实际应该使用PIL或其他图像库生成真实图片
        
        png_header = b'\x89PNG\r\n\x1a\n'
        placeholder_text = f"{text} - {width}x{height}".encode('utf-8')
        
        # 创建最小的PNG数据结构（这只是示例，不是完整的PNG）
        png_data = png_header + placeholder_text
        
        return png_data